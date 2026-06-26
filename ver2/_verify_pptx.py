# -*- coding: utf-8 -*-
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.stdout.reconfigure(encoding="utf-8")
p = Presentation(str(Path(__file__).parent / "report" / "ML Project.pptx"))
issues = []
for i, slide in enumerate(p.slides, 1):
    title = (slide.shapes.title.text if slide.shapes.title else "") or f"(slide {i})"
    texts = []
    npics = 0
    for sh in slide.shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t:
                texts.append(t)
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            npics += 1
    body = " | ".join(texts)
    for bad in ["23 meta", "24 inputs", "110M", "AI gen", "tui k bt", "thoại con", "not yet in the repo", "Next Step"]:
        if bad.lower() in body.lower():
            issues.append(f"Slide {i}: contains '{bad}'")
    print(f"{i:2d}. [{npics} pics] {title[:55]}")
    if i in (4, 13, 16, 19, 20, 21, 25, 26):
        for t in texts[:3]:
            print(f"     {t[:120]}")
print("\nIssues:", issues or "none")
