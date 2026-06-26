# -*- coding: utf-8 -*-
"""Final polish pass on ML Project.pptx."""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

sys.stdout.reconfigure(encoding="utf-8")

PPTX = Path(__file__).parent / "report" / "ML Project.pptx"
DL = Path(r"C:\Users\LOQ\Downloads\ML Project.pptx")
EDA = Path(__file__).parent / "eda" / "outputs" / "plots"

prs = Presentation(str(PPTX))

# Slide 4 — restore title, dedupe body
s4 = prs.slides[3]
if s4.shapes.title:
    s4.shapes.title.text = "DEFINITION & PURPOSE"
for sh in list(s4.shapes):
    if not sh.has_text_frame or sh == s4.shapes.title:
        continue
    t = sh.text_frame.text
    if "binary classification system that estimates" in t:
        sh._element.getparent().remove(sh._element)
    elif t.startswith("Definition: Binary classification"):
        sh.text_frame.text = (
            "Definition: Binary classification — depression-related (1) vs normal (0) "
            "from title+body + metadata (upvotes, comments, posting time).\n"
            "Mechanism: Two-stage pipeline — Stage 1 MentalRoBERTa+LoRA outputs p_text; "
            "Stage 2 LightGBM combines p_text + 38 metadata features (39 inputs total).\n"
            "Purpose: Subreddit-blind screening suitable for real-world posts beyond Reddit."
        )

# Slide 19 — Stage 2 feature set: add mutual-info plot (EDA context, not stage2 FI)
s19 = prs.slides[18]
if not any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in s19.shapes):
    s19.shapes.add_picture(str(EDA / "mutual_info_with_label.png"),
                           Inches(5.2), Inches(1.3), Inches(4.3), Inches(4.8))

# Slide 7 — data quality: add body length distribution
s7 = prs.slides[6]
if not any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in s7.shapes):
    s7.shapes.add_picture(str(EDA / "hist_body_chars.png"),
                          Inches(5.0), Inches(1.2), Inches(4.5), Inches(4.5))

# Slide 26 — demo: remove stale Reddit flow diagrams, keep confusion + text
s26 = prs.slides[25]
pics = [sh for sh in s26.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
for sh in pics[1:]:
    sh._element.getparent().remove(sh._element)

# Slide 25 — summary: simplify duplicate headers
s25 = prs.slides[24]
for sh in s25.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip() in ("Completed", "Future Work", "Open Orientation"):
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True

prs.save(str(PPTX))
import shutil
shutil.copy2(PPTX, DL)
print("Polished and saved.")
