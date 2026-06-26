# -*- coding: utf-8 -*-
"""Update ML Project.pptx — uses ORIGINAL slide indices, then reorders."""
from __future__ import annotations

import re
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

sys.stdout.reconfigure(encoding="utf-8")

ROOT = Path(__file__).resolve().parent
SRC = Path(r"C:\Users\LOQ\Downloads\ML Project.pptx")
BACKUP = Path(r"C:\Users\LOQ\Downloads\ML Project.backup.pptx")
OUT = ROOT / "report" / "ML Project.pptx"
FIG = ROOT / "report_figures"
EDA = ROOT / "eda" / "outputs" / "plots"
S1 = ROOT / "stage1" / "outputs" / "predictions"
S2R = ROOT / "stage2" / "outputs" / "report"

# Reorder AFTER edits. 0-based indices into ORIGINAL 27-slide deck.
NEW_ORDER = [
    0, 1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11,   # 1-12 EDA
    13, 14, 15, 20, 21, 22,                    # Stage 1 model + results
    12,                                        # Stage 2 feature set
    16, 17, 18, 19,                            # Stage 2 model + results
    23, 24, 25, 26,
]

REMOVE_PATTERNS = [
    r"AI gen v trminh check lại đúng lí do k nhe :\(",
    r"\(tui k bt vt j c",
    r"thoại con AI gợi ý đây mà Tuấn k bt đưa j vào, trminh cứu nhé\s*Ok?",
    r"Note: Stage 2 code is not yet in the repository[^\n]*",
    r"\(Optional\) Add a learning curve chart[^\n]*",
]


def title_of(slide) -> str:
    return (slide.shapes.title.text if slide.shapes.title else "") or ""


def reorder_slides(prs, order):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    lst.clear()
    for i in order:
        lst.append(ids[i])


def pics(slide):
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


def replace_pic(slide, idx, path: Path):
    p = pics(slide)
    if idx >= len(p) or not path.exists():
        print(f"  skip pic[{idx}] {path.name} on '{title_of(slide)[:35]}'")
        return
    old = p[idx]
    l, t, w, h = old.left, old.top, old.width, old.height
    old._element.getparent().remove(old._element)
    slide.shapes.add_picture(str(path), l, t, w, h)
    print(f"  pic[{idx}] <- {path.name}  ({title_of(slide)[:35]})")


def remove_pic(slide, idx):
    p = pics(slide)
    if idx < len(p):
        p[idx]._element.getparent().remove(p[idx]._element)
        print(f"  removed pic[{idx}] ({title_of(slide)[:35]})")


def set_text_in_slide(slide, contains: str, new: str):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if contains.lower() in sh.text_frame.text.lower():
            sh.text_frame.text = new
            return True
    return False


def clean_all_text(prs):
    repl = [
        ("23 metadata features", "38 metadata features (29 raw + 9 engineered)"),
        ("23 meta-features", "p_text + 38 metadata features (29 raw + 9 engineered)"),
        ("24 inputs", "39 inputs (p_text + 38 metadata)"),
        ("gets 24 inputs", "gets 39 inputs (p_text + 38 metadata)"),
        ("110M parameters", "~125M parameters"),
        ("110M", "~125M"),
    ]
    for slide in prs.slides:
        for sh in slide.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text
            for a, b in repl:
                t = t.replace(a, b)
            for pat in REMOVE_PATTERNS:
                t = re.sub(pat, "", t, flags=re.IGNORECASE)
            sh.text_frame.text = re.sub(r"\n{3,}", "\n\n", t).strip()


def main():
    OUT.parent.mkdir(parents=True, exist_ok=True)
    src = BACKUP if BACKUP.exists() else SRC
    shutil.copy2(src, OUT)
    prs = Presentation(str(OUT))
    s = prs.slides
    print(f"Source: {src.name} | slides: {len(s)}")

    clean_all_text(prs)

    # --- Text updates by original index ---
    set_text_in_slide(s[3], "Definition",
        "Definition: Binary classification — depression-related (1) vs normal (0) from title+body + metadata.\n"
        "Mechanism: Stage 1 MentalRoBERTa+LoRA → p_text; Stage 2 LightGBM on p_text + 38 metadata (39 inputs).\n"
        "Purpose: Subreddit-blind screening with modular text + meta ensemble.")

    set_text_in_slide(s[4], "complementary",
        "Text and metadata provide complementary signals\n"
        "Modular design — debug each stage independently\n"
        "Enables ablation: text-only vs meta-only vs full ensemble")

    set_text_in_slide(s[10], "RoBERTa",
        "RoBERTa fixed window — choose max_length balancing coverage vs compute.\n"
        "max_length=512 covers 96.6% of posts; head+tail truncation for the rest.")

    for sh in s[11].shapes:
        if sh.has_text_frame and ("Step" in sh.text_frame.text or "meta" in sh.text_frame.text.lower()):
            t = sh.text_frame.text
            t = re.sub(r"\b23\b", "38", t)
            t = re.sub(r"\b24\b inputs", "39 inputs", t)
            sh.text_frame.text = t

    for sh in s[12].shapes:
        if sh.has_text_frame and len(sh.text_frame.text) > 60:
            sh.text_frame.text = (
                "Stage 2 LightGBM: 39 inputs = p_text + 38 metadata (29 raw + 9 engineered).\n"
                "Rule: computable from one post alone — no subreddit, year, or subreddit-percentile proxies.\n"
                "Raw: length, engagement, lexical counts, temporal cyclical, booleans.\n"
                "Engineered: length-normalized rates (first_person_rate, negative_word_rate, …)."
            )

    set_text_in_slide(s[14], "Enhances generalization",
        "4-step head: layer-weighted avg (last 4 layers) → CLS + masked-mean pool → residual FFN → "
        "Multi-Sample Dropout (K=5).\n"
        "Enhances generalization without full fine-tuning of ~125M backbone parameters.")

    set_text_in_slide(s[16], "Input:",
        "Input: p_text (OOF) + 38 metadata (29 raw + 9 engineered) = 39 features\n"
        "Ablation: A=text only | B=meta only | C=full ensemble (production)\n"
        "Tuning: Optuna 50 trials, class_weight='balanced'\n"
        "Stacking on OOF p_text prevents leakage")

    set_text_in_slide(s[20], "Fold 0",
        "Per-fold @best-f1 (verified predictions):\n"
        "F0 OOF thr=0.69 f1=0.943 R_pos=0.910 | F0 TEST thr=0.62 f1=0.945 R_pos=0.925\n"
        "F1 OOF thr=0.77 f1=0.945 R_pos=0.918 | F1 TEST thr=0.73 f1=0.947 R_pos=0.929\n"
        "F2 OOF thr=0.87 f1=0.940 R_pos=0.906 | F2 TEST thr=0.84 f1=0.942 R_pos=0.915\n"
        "~6.4 h/fold Kaggle 2×T4 | Primary metric: F1-macro")

    set_text_in_slide(s[23], "Comparison",
        "Ablation TEST 2022 (subreddit-blind, thr=0.50):\n"
        "A text-only: f1_macro=0.940, f1_pos=0.914\n"
        "B meta-only: f1_macro=0.881, f1_pos=0.826\n"
        "C full ensemble: f1_macro=0.951, f1_pos=0.929, recall_pos=0.966")

    for sh in s[24].shapes:
        if not sh.has_text_frame:
            continue
        t = sh.text_frame.text
        if t.strip() == "Next Step":
            sh.text_frame.text = "Future Work"
        elif "Completed in-depth" in t:
            sh.text_frame.text = (
                "Completed: EDA + preprocessing; Stage 1 + Stage 2 trained & evaluated on TEST 2022; "
                "ablation A/B/C; live inference app (FB/Threads/Reddit, VI→EN, domain-mismatch handling)."
            )
        elif "Synthesize" in t or "LightGBM Stage 2" in t:
            sh.text_frame.text = (
                "Future: Platt/isotonic calibration, cross-lingual models, "
                "dataset balancing, cross-subreddit holdout."
            )

    # --- Images by original index ---
    replace_pic(s[5], 0, EDA / "top20_subreddits.png")
    replace_pic(s[7], 0, EDA / "label_rate_top30_subreddits.png")
    replace_pic(s[7], 1, EDA / "scatter_subreddit_n_vs_label_rate.png")
    replace_pic(s[8], 0, EDA / "distinctive_terms_label1_top25.png")
    replace_pic(s[8], 1, EDA / "lexical_lift_label1_over_label0.png")
    replace_pic(s[9], 0, EDA / "label_rate_by_year.png")
    replace_pic(s[9], 1, EDA / "posts_by_year.png")
    replace_pic(s[10], 0, EDA / "hist_token_lens_pair.png")

    remove_pic(s[13], 0)  # duplicate token histogram on Stage 1 intro

    remove_pic(s[16], 1)  # misleading EDA mutual-info on Stage 2 intro

    replace_pic(s[17], 0, S2R / "reliability_test.png")
    replace_pic(s[17], 1, FIG / "s2_confusion.png")
    box = s[17].shapes.add_textbox(Inches(0.4), Inches(5.65), Inches(9.2), Inches(1.25))
    box.text_frame.text = (
        "Calibration: baseline over-confident (mean_p≈0.29 vs base 0.26). "
        "Surveyed Platt/isotonic; kept baseline to prioritize positive recall."
    )
    for p in box.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(12)

    replace_pic(s[18], 0, FIG / "s2_feature_importance.png")
    replace_pic(s[19], 0, FIG / "s2_roc_pr.png")
    replace_pic(s[19], 1, FIG / "s2_score_dist.png")

    replace_pic(s[20], 0, S1 / "confusion_matrices.png")
    replace_pic(s[21], 0, FIG / "s1_confusion.png")
    replace_pic(s[21], 1, FIG / "s1_roc_pr.png")
    replace_pic(s[22], 0, FIG / "s1_calibration.png")
    replace_pic(s[22], 1, FIG / "s1_score_dist.png")
    replace_pic(s[23], 0, FIG / "ablation_abc.png")

    demo_box = s[25].shapes.add_textbox(Inches(0.35), Inches(0.85), Inches(9.3), Inches(2.5))
    demo_box.text_frame.text = (
        "Live inference app (ver2/app) — FastAPI + web UI\n"
        "• Scrape Reddit/Facebook/Threads URL or paste text + engagement metadata\n"
        "• Optional VI→EN translation → Stage 1 p_text → Stage 2 p_final\n"
        "• Non-Reddit posts: neutralize Reddit-specific meta features (domain-mismatch)\n"
        "• Production: full_trainable.safetensors + stage2_lgbm.pkl"
    )
    for p in demo_box.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(13)
    replace_pic(s[25], 0, FIG / "s2_confusion.png")

    print("Reordering slides...")
    reorder_slides(prs, NEW_ORDER)
    prs.save(str(OUT))
    shutil.copy2(OUT, SRC)

    prs2 = Presentation(str(OUT))
    print("\nFinal order:")
    for i, sl in enumerate(prs2.slides, 1):
        print(f"  {i:2d}. {title_of(sl)[:65]}")
    print(f"\nSaved: {OUT}")


if __name__ == "__main__":
    main()
