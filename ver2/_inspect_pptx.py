# -*- coding: utf-8 -*-
"""Inspect PPTX slide content for editing."""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.stdout.reconfigure(encoding="utf-8")

pptx_path = Path(r"C:\Users\LOQ\Downloads\ML Project.pptx")
p = Presentation(str(pptx_path))
print(f"slides: {len(p.slides)}\n")

for i, slide in enumerate(p.slides, 1):
    title = slide.shapes.title.text if slide.shapes.title else "(no title)"
    print(f"=== SLIDE {i}: {title} ===")
    for j, shape in enumerate(slide.shapes):
        st = shape.shape_type
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip().replace("\n", " | ")
            if txt:
                print(f"  text[{j}]: {txt[:300]}")
        if st == MSO_SHAPE_TYPE.PICTURE:
            print(f"  picture[{j}]: {shape.width}x{shape.height}")
    print()
